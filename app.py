import streamlit as st
import os
from pypdf import PdfReader
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import groq
from dotenv import load_dotenv
import tempfile
import pickle
import requests
from bs4 import BeautifulSoup
import urllib.parse
from typing import List, Tuple, Optional, Union
import docx
import pptx
import time
import threading
from queue import Queue
from urllib.parse import urljoin, urlparse
import concurrent.futures
import io

# Load environment variables
load_dotenv()

@st.cache_resource
def get_embedding_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

class DocumentProcessor:
    def __init__(self):
        self.client = groq.Groq(api_key=os.getenv("GROQ_API_KEY"))
        self.vector_dimension = 384
        self.embedding_model = get_embedding_model()
        self.index = faiss.IndexFlatL2(self.vector_dimension)
        self.text_chunks = []
        self.chunk_sources = []
        self.processing_queue = Queue()
        self.processed_urls = set()

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        try:
            # Use BytesIO to handle PDF content in memory
            with io.BytesIO(file_content) as pdf_file:
                pdf_reader = PdfReader(pdf_file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            st.error(f"Error processing PDF: {str(e)}")
            return ""

    def extract_text_from_docx(self, file_content: bytes) -> str:
        try:
            # Use BytesIO to handle DOCX content in memory
            with io.BytesIO(file_content) as docx_file:
                doc = docx.Document(docx_file)
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            st.error(f"Error processing DOCX: {str(e)}")
            return ""

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        try:
            # Use BytesIO to handle PPTX content in memory
            with io.BytesIO(file_content) as pptx_file:
                prs = pptx.Presentation(pptx_file)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            return text
        except Exception as e:
            st.error(f"Error processing PPTX: {str(e)}")
            return ""

    def extract_text_from_txt(self, file_content: bytes) -> str:
        try:
            return file_content.decode('utf-8')
        except Exception as e:
            st.error(f"Error processing TXT: {str(e)}")
            return ""

    def process_file(self, file_content: bytes, file_name: str) -> str:
        """Process different file types based on extension."""
        extension = file_name.lower().split('.')[-1]
        
        extractors = {
            'pdf': self.extract_text_from_pdf,
            'docx': self.extract_text_from_docx,
            'pptx': self.extract_text_from_pptx,
            'txt': self.extract_text_from_txt
        }
        
        extractor = extractors.get(extension)
        if extractor:
            return extractor(file_content)
        else:
            st.error(f"Unsupported file type: {extension}")
            return ""

    def split_text(self, text: str, chunk_size: int = 1000) -> List[str]:
        if not text.strip():
            return []
        words = text.split()
        chunks = []
        current_chunk = []
        current_size = 0
        
        for word in words:
            current_chunk.append(word)
            current_size += len(word) + 1
            
            if current_size >= chunk_size:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_size = 0
        
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        return chunks

    def delayed_process_documents(self):
        """Process documents from queue with delay."""
        while True:
            if not self.processing_queue.empty():
                content, source_name = self.processing_queue.get()
                time.sleep(30)  # 30-second delay
                self.process_document(content, source_name)
                self.save_index()
            time.sleep(1)

    def start_processing_thread(self):
        """Start the background processing thread."""
        thread = threading.Thread(target=self.delayed_process_documents, daemon=True)
        thread.start()

    def process_document(self, content: str, source_name: str) -> int:
        try:
            chunks = self.split_text(content)
            
            if not chunks:
                return 0
                
            embeddings = self.embedding_model.encode(chunks)
            faiss.normalize_L2(embeddings)
            
            self.index.add(embeddings.astype('float32'))
            
            self.text_chunks.extend(chunks)
            self.chunk_sources.extend([source_name] * len(chunks))
            
            return len(chunks)
        except Exception as e:
            st.error(f"Error processing document {source_name}: {str(e)}")
            return 0

    def save_index(self, directory: str = "./saved_index") -> None:
        try:
            os.makedirs(directory, exist_ok=True)
            faiss.write_index(self.index, os.path.join(directory, "docs.index"))
            with open(os.path.join(directory, "chunks.pkl"), "wb") as f:
                pickle.dump((self.text_chunks, self.chunk_sources), f)
        except Exception as e:
            st.error(f"Error saving index: {str(e)}")

    def load_index(self, directory: str = "./saved_index") -> bool:
        try:
            index_path = os.path.join(directory, "docs.index")
            chunks_path = os.path.join(directory, "chunks.pkl")
            
            if os.path.exists(index_path) and os.path.exists(chunks_path):
                self.index = faiss.read_index(index_path)
                with open(chunks_path, "rb") as f:
                    self.text_chunks, self.chunk_sources = pickle.load(f)
                return True
            return False
        except Exception as e:
            st.error(f"Error loading index: {str(e)}")
            return False

    def query_documents(self, query: str, n_results: int = 3) -> Tuple[str, List[Tuple[str, str]]]:
        try:
            query_embedding = self.embedding_model.encode([query])
            faiss.normalize_L2(query_embedding)
            
            distances, indices = self.index.search(query_embedding.astype('float32'), n_results)
            relevant_chunks = [self.text_chunks[i] for i in indices[0]]
            chunk_sources_info = [self.chunk_sources[i] for i in indices[0]]
            
            context = "\n\n".join(relevant_chunks)
            prompt = f"""You are a helpful assistant to help users with their questions. 
            Based on the following context, answer the question: {query}

            Context:
            {context}

            Answer:"""
            
            completion = self.client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.1-8b-instant",
                temperature=0.8,
            )
            
            return completion.choices[0].message.content, list(zip(relevant_chunks, chunk_sources_info))
        except Exception as e:
            st.error(f"Error querying documents: {str(e)}")
            return "I encountered an error while processing your question. Please try again.", []

    def crawl_website(self, base_url: str, max_pages: int = 50) -> List[Tuple[str, str]]:
        """Crawl website and extract text from pages."""
        pages_content = []
        to_visit = {base_url}
        visited = set()

        while to_visit and len(visited) < max_pages:
            url = to_visit.pop()
            if url in visited:
                continue

            try:
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                }
                response = requests.get(url, headers=headers, timeout=10)
                response.raise_for_status()
                
                soup = BeautifulSoup(response.text, 'html.parser')
                for script in soup(["script", "style", "meta", "noscript"]):
                    script.decompose()
                text = soup.get_text()
                pages_content.append((url, text))

                # Extract and add new links
                links = soup.find_all('a', href=True)
                for link in links:
                    href = urljoin(url, link['href'])
                    if self.is_same_domain(url, href) and href not in visited:
                        to_visit.add(href)
                
                visited.add(url)
                time.sleep(1)  # Be nice to servers
                
            except Exception as e:
                st.error(f"Error crawling {url}: {str(e)}")
                continue

        return pages_content

    def is_same_domain(self, url1: str, url2: str) -> bool:
        """Check if two URLs belong to the same domain."""
        domain1 = urlparse(url1).netloc
        domain2 = urlparse(url2).netloc
        return domain1 == domain2

def main():
    st.set_page_config(page_title="Groq RAG", page_icon="🎓", layout="wide")
    
    st.title("Realtime RAG")
    
    if 'processor' not in st.session_state:
        st.session_state.processor = DocumentProcessor()
        st.session_state.processor.start_processing_thread()
    
    if not st.session_state.get("index_loaded"):
        if st.session_state.processor.load_index():
            st.session_state.index_loaded = True
            st.success("Loaded existing knowledge base!")
        else:
            st.info("No existing knowledge base found. Please upload documents or add URLs.")

    with st.sidebar:
        st.header("Add Documents")
        
        supported_types = ["pdf", "docx", "txt", "pptx"]
        uploaded_files = st.file_uploader(
            f"Upload documents ({', '.join(supported_types)})", 
            type=supported_types, 
            accept_multiple_files=True
        )
        
        st.header("Add URLs")
        url_input = st.text_area("Enter URLs (one per line)")
        
        crawl_depth = st.slider("Maximum pages to crawl per website", 1, 100, 50)
        
        process_button = st.button("Process Documents")
        
        if process_button:
            total_queued = 0
            
            # Process uploaded files
            if uploaded_files:
                for uploaded_file in uploaded_files:
                    with st.spinner(f"Queueing {uploaded_file.name}..."):
                        # Read file content directly into memory
                        file_content = uploaded_file.read()
                        text = st.session_state.processor.process_file(file_content, uploaded_file.name)
                        if text:
                            st.session_state.processor.processing_queue.put((text, uploaded_file.name))
                            total_queued += 1
            
            # Process URLs
            if url_input:
                urls = [url.strip() for url in url_input.split('\n') if url.strip()]
                for url in urls:
                    if url:
                        with st.spinner(f"Crawling {url}..."):
                            pages_content = st.session_state.processor.crawl_website(url, max_pages=crawl_depth)
                            for page_url, content in pages_content:
                                if content:
                                    st.session_state.processor.processing_queue.put((content, page_url))
                                    total_queued += 1
            
            if total_queued > 0:
                st.success(f"Queued {total_queued} documents for processing. They will be processed with a 30-second delay.")
            else:
                st.warning("No new content was queued. Please check your inputs.")

        if st.button("Clear Knowledge Base"):
            if os.path.exists("./saved_index"):
                try:
                    os.remove("./saved_index/docs.index")
                    os.remove("./saved_index/chunks.pkl")
                    st.session_state.processor = DocumentProcessor()
                    st.session_state.processor.start_processing_thread()
                    st.session_state.index_loaded = False
                    st.success("Knowledge base cleared successfully!")
                except Exception as e:
                    st.error(f"Error clearing knowledge base: {str(e)}")

    st.header("Ask Questions")
    query = st.text_input("Enter your question:")

    if query:
        if len(st.session_state.processor.text_chunks) == 0:
            st.warning("Please add some documents to the knowledge base first!")
        else:
            with st.spinner("Searching and generating response..."):
                answer, sources = st.session_state.processor.query_documents(query)
                
                st.markdown("### Answer")
                st.write(answer)
                
                st.markdown("### Sources")
                for i, (chunk, source) in enumerate(sources, 1):
                    with st.expander(f"Source {i} - {source}"):
                        st.write(chunk)

if __name__ == "__main__":
    main()
